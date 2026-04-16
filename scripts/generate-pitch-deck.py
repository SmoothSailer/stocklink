"""Generate Ristoka investor pitch deck as .pptx"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Brand colors ──
GREEN       = RGBColor(0x16, 0xA3, 0x4A)
GREEN_DARK  = RGBColor(0x15, 0x80, 0x3D)
GREEN_LIGHT = RGBColor(0xDC, 0xFC, 0xE7)
ORANGE      = RGBColor(0xF9, 0x73, 0x16)
RED         = RGBColor(0xEF, 0x44, 0x44)
BLUE        = RGBColor(0x3B, 0x82, 0xF6)
GRAY_50     = RGBColor(0xF9, 0xFA, 0xFB)
GRAY_100    = RGBColor(0xF3, 0xF4, 0xF6)
GRAY_200    = RGBColor(0xE5, 0xE7, 0xEB)
GRAY_500    = RGBColor(0x6B, 0x72, 0x80)
GRAY_700    = RGBColor(0x37, 0x41, 0x51)
GRAY_800    = RGBColor(0x1F, 0x29, 0x37)
GRAY_900    = RGBColor(0x11, 0x18, 0x27)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]  # Blank

# ── Helpers ──

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_bg(slide, c1, c2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = c1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2
    fill.gradient_stops[1].position = 1.0

def add_text_box(slide, left, top, width, height, text, font_size=18, color=GRAY_800, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_multiline_box(slide, left, top, width, height, lines, font_size=16, color=GRAY_700, line_spacing=1.3, bold_first=False, font_name="Calibri"):
    """lines = list of strings"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox

def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Smaller corner radius
    shape.adjustments[0] = 0.05
    return shape

def add_circle(slide, left, top, size, fill_color):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape

def add_slide_number(slide, num, total, color=GRAY_500):
    add_text_box(slide, Inches(12.0), Inches(7.0), Inches(1.2), Inches(0.4),
                 f"{num} / {total}", font_size=11, color=color, alignment=PP_ALIGN.RIGHT)

def add_brand_tag(slide, color=GRAY_500):
    add_text_box(slide, Inches(0.5), Inches(7.0), Inches(1.5), Inches(0.4),
                 "RISTOKA", font_size=11, color=color, bold=True)

def add_section_label(slide, text, color=GREEN):
    add_text_box(slide, Inches(0.8), Inches(0.6), Inches(4), Inches(0.4),
                 text.upper(), font_size=12, color=color, bold=True)

TOTAL_SLIDES = 14

# ══════════════════════════════════════════════
# SLIDE 1: TITLE
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_gradient_bg(slide, GREEN, GREEN_DARK)

# Logo mark
logo = add_rounded_rect(slide, Inches(0.8), Inches(1.5), Inches(0.9), Inches(0.9), RGBColor(0xFF, 0xFF, 0xFF))
logo.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
# R on the logo - use a separate transparent overlay
add_text_box(slide, Inches(0.8), Inches(1.48), Inches(0.9), Inches(0.9),
             "R", font_size=40, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(2.8), Inches(8), Inches(1.2),
             "Ristoka", font_size=56, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(4.0), Inches(8), Inches(0.7),
             "Stocking you up, backing you up.", font_size=28, color=RGBColor(0xDC, 0xFC, 0xE7))
add_text_box(slide, Inches(0.8), Inches(4.9), Inches(8), Inches(0.5),
             "The supply + financing platform for African retailers", font_size=18, color=RGBColor(0xBB, 0xF7, 0xD0))
add_slide_number(slide, 1, TOTAL_SLIDES, RGBColor(0xBB, 0xF7, 0xD0))
add_brand_tag(slide, RGBColor(0xBB, 0xF7, 0xD0))

# ══════════════════════════════════════════════
# SLIDE 2: THE PROBLEM
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide, GRAY_50)
add_section_label(slide, "The Problem")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "Retailers are stuck in a cycle of scarcity", font_size=40, color=GRAY_900, bold=True)

# Three problem cards
cards = [
    ("💰", "No Working Capital", "Retailers buy 2 bags when they could sell 10.\nBanks won't lend. M-Shwari caps at KSh 5,000.\nThey need KSh 30,000–50,000/week."),
    ("🔗", "Middleman Tax", "No direct wholesaler access.\nEvery layer adds 10-20% markup.\nA KSh 3,200 bag of rice costs her KSh 3,800."),
    ("🌍", "No Global Access", "International products are completely\ninaccessible. No logistics, no contacts,\nno capital to import."),
]
for i, (icon, title, desc) in enumerate(cards):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.6)
    # Card bg
    add_rounded_rect(slide, x, y, Inches(3.6), Inches(3.4), WHITE, GRAY_200)
    # Red left border
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.3), Inches(0.06), Inches(2.8)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = RED
    slide.shapes[-1].line.fill.background()
    # Icon
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.5),
                 icon, font_size=28, color=GRAY_800)
    # Title
    add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.0), Inches(0.4),
                 title, font_size=18, color=GRAY_900, bold=True)
    # Desc
    add_multiline_box(slide, x + Inches(0.3), y + Inches(1.4), Inches(3.0), Inches(2.0),
                      desc.split("\n"), font_size=14, color=GRAY_500)

add_slide_number(slide, 2, TOTAL_SLIDES)
add_brand_tag(slide)

# ══════════════════════════════════════════════
# SLIDE 3: THE OPPORTUNITY
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide, GRAY_900)
add_section_label(slide, "The Opportunity", GREEN)

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "A massive, underserved market", font_size=40, color=WHITE, bold=True)

stats = [
    ("$200B+", "Africa's informal retail market"),
    ("80%", "Of retail is informal / small shops"),
    ("~$0", "Credit available to micro-retailers"),
]
for i, (val, label) in enumerate(stats):
    x = Inches(0.8 + i * 4.0)
    add_rounded_rect(slide, x, Inches(2.8), Inches(3.6), Inches(2.2),
                     RGBColor(0x1E, 0x29, 0x3B), RGBColor(0x37, 0x41, 0x51))
    add_text_box(slide, x + Inches(0.4), Inches(3.1), Inches(2.8), Inches(0.8),
                 val, font_size=42, color=RGBColor(0x4A, 0xDE, 0x80), bold=True)
    add_text_box(slide, x + Inches(0.4), Inches(3.9), Inches(2.8), Inches(0.6),
                 label, font_size=15, color=GRAY_500)

add_text_box(slide, Inches(0.8), Inches(5.6), Inches(9), Inches(0.5),
             "Every retailer who can stock more, sells more. The constraint isn't demand — it's supply and capital.",
             font_size=17, color=GRAY_500)

add_slide_number(slide, 3, TOTAL_SLIDES, GRAY_500)
add_brand_tag(slide, GRAY_500)

# ══════════════════════════════════════════════
# SLIDE 4: OUR SOLUTION
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide, WHITE)
add_section_label(slide, "Our Solution")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "Supply + Capital in one platform", font_size=40, color=GRAY_900, bold=True)

solutions = [
    ("📦", "Supply Platform", "Aggregated catalog from local\nwholesalers, manufacturers, and\ninternational suppliers. Real prices,\nreal stock, real-time."),
    ("🏦", "Embedded Financing", "Murabaha-based BNPL — we buy\ninventory and sell at a transparent\nmarkup. Not interest. 30% down,\nrest in installments."),
    ("🌐", "Global Sourcing", "Products from UK, Turkey, UAE\nwith logistics handled. Retailers\naccess international goods they\ncouldn't reach on their own."),
]
for i, (icon, title, desc) in enumerate(solutions):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.6)
    add_rounded_rect(slide, x, y, Inches(3.6), Inches(3.6), WHITE, GRAY_200)
    # Green left border
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.3), Inches(0.06), Inches(3.0)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = GREEN
    slide.shapes[-1].line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.6), Inches(0.5),
                 icon, font_size=28, color=GRAY_800)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.9), Inches(3.0), Inches(0.4),
                 title, font_size=18, color=GRAY_900, bold=True)
    add_multiline_box(slide, x + Inches(0.3), y + Inches(1.45), Inches(3.0), Inches(2.2),
                      desc.split("\n"), font_size=14, color=GRAY_500)

add_slide_number(slide, 4, TOTAL_SLIDES)
add_brand_tag(slide)

# ══════════════════════════════════════════════
# SLIDE 5: HOW IT WORKS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide, GRAY_50)
add_section_label(slide, "How It Works")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "Simple flow, real transactions", font_size=40, color=GRAY_900, bold=True)

steps = [
    ("1", "Browse & Order", "Retailer searches products,\nselects quantity & unit,\nconfirms order"),
    ("2", "Choose Payment", "M-Pesa, Cash, Card,\nor BNPL (Buy Now\nPay Later)"),
    ("3", "WhatsApp Confirm", "Order goes to assigned\nsales rep via WhatsApp\nfor fulfillment"),
    ("4", "Track & Pay", "Real-time order tracking.\nBNPL installments\nmanaged in-app"),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.6 + i * 3.15)
    # Circle with number
    circle = add_circle(slide, x + Inches(1.05), Inches(2.8), Inches(0.55), GREEN)
    add_text_box(slide, x + Inches(1.05), Inches(2.82), Inches(0.55), Inches(0.55),
                 num, font_size=22, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.2), Inches(3.5), Inches(2.7), Inches(0.4),
                 title, font_size=17, color=GRAY_900, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline_box(slide, x + Inches(0.2), Inches(4.0), Inches(2.7), Inches(1.5),
                      desc.split("\n"), font_size=13, color=GRAY_500)
    for p in slide.shapes[-1].text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER
    # Arrow between steps
    if i < 3:
        add_text_box(slide, x + Inches(2.85), Inches(2.9), Inches(0.4), Inches(0.4),
                     "→", font_size=24, color=GRAY_200, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(6.0), Inches(11.5), Inches(0.4),
             "WhatsApp-first UX — meets retailers where they already are",
             font_size=15, color=GRAY_500, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 5, TOTAL_SLIDES)
add_brand_tag(slide)

# ══════════════════════════════════════════════
# SLIDE 6: PRODUCT IS BUILT
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide, WHITE)
add_section_label(slide, "Product")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "This is live, not a prototype", font_size=40, color=GRAY_900, bold=True)

features = [
    "✓  Full product catalog with categories, search, filters",
    "✓  Multi-unit pricing (pieces, boxes, cartons)",
    "✓  Flash deals with countdown timers",
    "✓  Order tracking with real-time status updates",
    "✓  BNPL with KYC verification & credit limits",
    "✓  Admin dashboard — orders, inventory, retailers",
    "✓  Sales rep portal — onboarding, payments, BNPL",
    "✓  Wholesaler read-only inventory view",
    "✓  Affiliate referral system",
]
add_multiline_box(slide, Inches(0.8), Inches(2.5), Inches(7), Inches(4.5),
                  features, font_size=16, color=GRAY_700, line_spacing=1.6)

# "Live Demo" card on right
add_rounded_rect(slide, Inches(9.0), Inches(2.5), Inches(3.5), Inches(3.5), GRAY_100, GRAY_200)
add_rounded_rect(slide, Inches(9.3), Inches(2.8), Inches(2.9), Inches(1.8), GREEN)
add_text_box(slide, Inches(9.3), Inches(3.0), Inches(2.9), Inches(0.8),
             "R", font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.3), Inches(3.8), Inches(2.9), Inches(0.4),
             "Live Demo Available", font_size=13, color=RGBColor(0xBB, 0xF7, 0xD0), alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(9.3), Inches(4.9), Inches(2.9), Inches(0.6),
             "Retailer app • Admin dashboard\n• Sales rep portal", font_size=12, color=GRAY_500, alignment=PP_ALIGN.CENTER)

# Tech footnote
add_text_box(slide, Inches(0.8), Inches(6.2), Inches(7), Inches(0.5),
             "Built with Next.js, Supabase, TypeScript. Mobile-first. 33 routes.",
             font_size=13, color=GRAY_500)

add_slide_number(slide, 6, TOTAL_SLIDES)
add_brand_tag(slide)

# ══════════════════════════════════════════════
# SLIDE 7: BNPL / MURABAHA
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide, GREEN_DARK)
add_section_label(slide, "BNPL Model", RGBColor(0xBB, 0xF7, 0xD0))

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "Murabaha — Transparent, Ethical Financing", font_size=38, color=WHITE, bold=True)

bnpl_steps = [
    ("🔍", "KYC Verification", "ID upload, business docs,\nadmin review. Credit limit\nset per retailer."),
    ("📦", "We Buy Inventory", "Ristoka purchases stock at\nwholesale cost on behalf\nof the retailer."),
    ("💵", "30% Down Payment", "Retailer pays 30% upfront.\nReduces our risk exposure\nimmediately."),
    ("📅", "Installments", "Remaining 70% + markup\nspread over 2-4 weekly\ninstallments."),
]
for i, (icon, title, desc) in enumerate(bnpl_steps):
    x = Inches(0.5 + i * 3.15)
    add_rounded_rect(slide, x, Inches(2.6), Inches(2.9), Inches(2.8),
                     RGBColor(0x12, 0x6B, 0x32), RGBColor(0x1A, 0x8B, 0x45))
    add_text_box(slide, x + Inches(0.3), Inches(2.8), Inches(2.3), Inches(0.5),
                 icon, font_size=30, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), Inches(3.3), Inches(2.3), Inches(0.4),
                 title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_multiline_box(slide, x + Inches(0.3), Inches(3.8), Inches(2.3), Inches(1.5),
                      desc.split("\n"), font_size=13, color=RGBColor(0xBB, 0xF7, 0xD0))
    for p in slide.shapes[-1].text_frame.paragraphs:
        p.alignment = PP_ALIGN.CENTER

add_text_box(slide, Inches(0.8), Inches(5.9), Inches(10), Inches(0.8),
             "This is cost-plus, not interest. The markup is transparent and agreed upfront.\nShariah-compliant. Trust-building. Retailer-friendly.",
             font_size=17, color=RGBColor(0xDC, 0xFC, 0xE7))

add_slide_number(slide, 7, TOTAL_SLIDES, RGBColor(0xBB, 0xF7, 0xD0))
add_brand_tag(slide, RGBColor(0xBB, 0xF7, 0xD0))

# ══════════════════════════════════════════════
# SLIDE 8: UNIT ECONOMICS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide, WHITE)
add_section_label(slide, "Unit Economics")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "Every transaction generates revenue", font_size=40, color=GRAY_900, bold=True)

# Table
rows_data = [
    ("", "Cash Order", "BNPL Order"),
    ("Order Value (10 bags rice)", "KSh 35,000", "KSh 35,000"),
    ("Our Cost (wholesale)", "KSh 32,000", "KSh 32,000"),
    ("Supply Margin (8-15%)", "KSh 3,000", "KSh 3,000"),
    ("BNPL Markup (10-15%)", "—", "KSh 3,500"),
    ("Total Revenue", "KSh 3,000", "KSh 6,500"),
    ("Return on Capital", "9.4%", "20.3%"),
    ("Cycle Time", "Immediate", "2-4 weeks"),
]

table_shape = slide.shapes.add_table(len(rows_data), 3, Inches(1.5), Inches(2.4), Inches(10), Inches(4.0))
table = table_shape.table

for r, row in enumerate(rows_data):
    for c, val in enumerate(row):
        cell = table.cell(r, c)
        cell.text = val
        p = cell.text_frame.paragraphs[0]
        p.font.size = Pt(15)
        p.font.name = "Calibri"
        if r == 0:  # Header
            p.font.bold = True
            p.font.color.rgb = GRAY_700
            cell.fill.solid()
            cell.fill.fore_color.rgb = GRAY_100
        elif r in (5, 6):  # Total / Return rows
            p.font.bold = True
            if c > 0:
                p.font.color.rgb = GREEN
            else:
                p.font.color.rgb = GRAY_900
        elif r == 3 or r == 4:  # Margin rows
            if c > 0 and val != "—":
                p.font.color.rgb = GREEN
                p.font.bold = True
            else:
                p.font.color.rgb = GRAY_700
        else:
            p.font.color.rgb = GRAY_700

add_text_box(slide, Inches(1.5), Inches(6.6), Inches(10), Inches(0.4),
             "Capital deployed into inventory generates 9-20% returns per cycle. Cycles repeat every 2-4 weeks.",
             font_size=14, color=GRAY_500)

add_slide_number(slide, 8, TOTAL_SLIDES)
add_brand_tag(slide)

# ══════════════════════════════════════════════
# SLIDE 9: REVENUE MODEL
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide, GRAY_50)
add_section_label(slide, "Revenue Model")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "Three revenue streams, all transactional", font_size=40, color=GRAY_900, bold=True)

streams = [
    ("①", "Supply Margin", "Buy wholesale, sell to retailers\nat 8-15% markup. Traditional\ndistribution economics."),
    ("②", "BNPL Markup", "Murabaha markup of 10-15% on\nfinanced orders. Revenue\nrecognized on delivery."),
    ("③", "Global Sourcing", "Higher margins on international\nproducts (20-40%). UK, Turkey,\nUAE sourcing."),
]
for i, (num, title, desc) in enumerate(streams):
    x = Inches(0.8 + i * 4.0)
    add_text_box(slide, x, Inches(2.6), Inches(0.8), Inches(0.8),
                 num, font_size=44, color=GREEN, bold=True)
    add_text_box(slide, x, Inches(3.4), Inches(3.5), Inches(0.4),
                 title, font_size=20, color=GRAY_900, bold=True)
    add_multiline_box(slide, x, Inches(3.9), Inches(3.5), Inches(1.5),
                      desc.split("\n"), font_size=15, color=GRAY_500)

# Insight callout
add_rounded_rect(slide, Inches(2.0), Inches(5.8), Inches(9.0), Inches(0.8), GREEN_LIGHT)
add_text_box(slide, Inches(2.3), Inches(5.9), Inches(8.5), Inches(0.6),
             "💡 Key insight: The capital itself is productive. Every shilling deployed generates margin revenue.",
             font_size=15, color=GREEN_DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 9, TOTAL_SLIDES)
add_brand_tag(slide)

# ══════════════════════════════════════════════
# SLIDE 10: RISK MITIGATION
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide, WHITE)
add_section_label(slide, "Risk Mitigation")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "Built-in safeguards at every layer", font_size=40, color=GRAY_900, bold=True)

safeguards = [
    "✓  30% mandatory down payment — immediate capital recovery",
    "✓  KYC verification — ID, business docs, admin review",
    "✓  Individual credit limits — set per retailer, tracked real-time",
    "✓  Overdue detection — blocked from new BNPL if overdue",
    "✓  Inventory as collateral — liquidatable physical goods",
    "✓  Short cycles — 2-4 week repayment, not months",
]
add_multiline_box(slide, Inches(0.8), Inches(2.5), Inches(7.5), Inches(3.5),
                  safeguards, font_size=16, color=GRAY_700, line_spacing=1.7)

# Default scenario card
add_rounded_rect(slide, Inches(9.0), Inches(2.5), Inches(3.8), Inches(3.8), GRAY_50, GRAY_200)
add_text_box(slide, Inches(9.3), Inches(2.7), Inches(3.2), Inches(0.4),
             "Default Scenario", font_size=16, color=GRAY_800, bold=True)
add_text_box(slide, Inches(9.3), Inches(3.2), Inches(3.2), Inches(0.5),
             "If retailer defaults after 30% down\non KSh 35,000 order:", font_size=13, color=GRAY_500)

scenario_lines = [
    ("Down payment collected", "KSh 10,500", GREEN),
    ("Max exposure", "KSh 24,500", ORANGE),
    ("Our cost was", "KSh 32,000", GRAY_700),
]
for j, (label, val, color) in enumerate(scenario_lines):
    y = Inches(4.0 + j * 0.45)
    add_text_box(slide, Inches(9.3), y, Inches(2.0), Inches(0.35),
                 label, font_size=12, color=GRAY_500)
    add_text_box(slide, Inches(11.3), y, Inches(1.4), Inches(0.35),
                 val, font_size=12, color=color, bold=True, alignment=PP_ALIGN.RIGHT)

add_text_box(slide, Inches(9.3), Inches(5.5), Inches(3.2), Inches(0.5),
             "Risk buffer (5-10% of fund) covers this.", font_size=12, color=GRAY_500)

add_slide_number(slide, 10, TOTAL_SLIDES)
add_brand_tag(slide)

# ══════════════════════════════════════════════
# SLIDE 11: THE ASK
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_gradient_bg(slide, GREEN, GREEN_DARK)
add_section_label(slide, "The Ask", RGBColor(0xBB, 0xF7, 0xD0))

add_text_box(slide, Inches(0.8), Inches(1.6), Inches(11), Inches(1.0),
             "$30,000 – $60,000", font_size=56, color=WHITE, bold=True)
add_text_box(slide, Inches(0.8), Inches(2.8), Inches(8), Inches(0.5),
             "Seed investment for 20% equity", font_size=24, color=RGBColor(0xDC, 0xFC, 0xE7))

# Option A
add_rounded_rect(slide, Inches(0.8), Inches(3.8), Inches(5.0), Inches(2.5),
                 RGBColor(0x12, 0x6B, 0x32), RGBColor(0x1A, 0x8B, 0x45))
add_text_box(slide, Inches(1.1), Inches(3.95), Inches(4.4), Inches(0.3),
             "STRUCTURE OPTION A", font_size=11, color=RGBColor(0xBB, 0xF7, 0xD0))
add_text_box(slide, Inches(1.1), Inches(4.3), Inches(4.4), Inches(0.4),
             "Pure Equity", font_size=22, color=WHITE, bold=True)
add_text_box(slide, Inches(1.1), Inches(4.85), Inches(4.4), Inches(0.8),
             "20% ownership. Returns via company\ngrowth and future liquidity events.",
             font_size=15, color=RGBColor(0xBB, 0xF7, 0xD0))

# Option B
add_rounded_rect(slide, Inches(6.3), Inches(3.8), Inches(5.5), Inches(2.5),
                 RGBColor(0x12, 0x6B, 0x32), RGBColor(0x1A, 0x8B, 0x45))
add_text_box(slide, Inches(6.6), Inches(3.95), Inches(4.9), Inches(0.3),
             "STRUCTURE OPTION B", font_size=11, color=RGBColor(0xBB, 0xF7, 0xD0))
add_text_box(slide, Inches(6.6), Inches(4.3), Inches(4.9), Inches(0.4),
             "Hybrid", font_size=22, color=WHITE, bold=True)
add_text_box(slide, Inches(6.6), Inches(4.85), Inches(4.9), Inches(0.8),
             "Equity + profit-share on BNPL revenue until\n1.5x capital returned. Then equity upside continues.",
             font_size=15, color=RGBColor(0xBB, 0xF7, 0xD0))

add_slide_number(slide, 11, TOTAL_SLIDES, RGBColor(0xBB, 0xF7, 0xD0))
add_brand_tag(slide, RGBColor(0xBB, 0xF7, 0xD0))

# ══════════════════════════════════════════════
# SLIDE 12: USE OF FUNDS
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide, WHITE)
add_section_label(slide, "Use of Funds")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "Capital goes where it generates returns", font_size=40, color=GRAY_900, bold=True)

# Stacked bar
bar_items = [
    (0.65, GREEN, "Inventory 65%"),
    (0.12, BLUE, "Logistics 12%"),
    (0.12, ORANGE, "Sales 12%"),
    (0.06, GRAY_700, "Tech 6%"),
    (0.05, RED, "Buffer 5%"),
]
bar_x = Inches(0.8)
bar_y = Inches(2.5)
bar_w_total = Inches(11.5)
bar_h = Inches(0.55)
for pct, color, label in bar_items:
    w = int(bar_w_total * pct)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, bar_x, bar_y, w, bar_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.fill.background()
    # Label inside bar
    if pct >= 0.10:
        add_text_box(slide, bar_x + Inches(0.1), bar_y + Inches(0.05), w - Inches(0.2), bar_h,
                     label, font_size=11, color=WHITE, bold=True)
    bar_x += w

# Breakdown grid
fund_items = [
    (GREEN, "Inventory", "Purchase stock for retailer\norders & BNPL financing"),
    (BLUE, "Logistics", "Delivery infrastructure\n& Sambaza integration"),
    (ORANGE, "Sales", "Sales reps on ground,\nretailer acquisition"),
    (GRAY_700, "Technology", "Hosting, M-Pesa STK Push,\npayment infrastructure"),
    (RED, "Risk Buffer", "Default protection,\nunexpected costs"),
]
for i, (color, title, desc) in enumerate(fund_items):
    x = Inches(0.8 + i * 2.4)
    # Color dot
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, Inches(3.6), Inches(0.18), Inches(0.18))
    dot.fill.solid()
    dot.fill.fore_color.rgb = color
    dot.line.fill.background()
    add_text_box(slide, x, Inches(3.9), Inches(2.2), Inches(0.3),
                 title, font_size=14, color=GRAY_800, bold=True)
    add_multiline_box(slide, x, Inches(4.3), Inches(2.2), Inches(0.8),
                      desc.split("\n"), font_size=12, color=GRAY_500)

# Callout
add_rounded_rect(slide, Inches(2.0), Inches(5.8), Inches(9.0), Inches(0.8), GREEN_LIGHT)
add_text_box(slide, Inches(2.3), Inches(5.9), Inches(8.5), Inches(0.6),
             "This is not a cash-burn startup. 65% of investment goes into revenue-generating inventory from day one.",
             font_size=15, color=GREEN_DARK, bold=True, alignment=PP_ALIGN.CENTER)

add_slide_number(slide, 12, TOTAL_SLIDES)
add_brand_tag(slide)

# ══════════════════════════════════════════════
# SLIDE 13: ROADMAP
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg(slide, WHITE)
add_section_label(slide, "Roadmap")

add_text_box(slide, Inches(0.8), Inches(1.2), Inches(11), Inches(0.8),
             "90-day plan post-investment", font_size=40, color=GRAY_900, bold=True)

months = [
    (GREEN, "MONTH 1", "Launch Operations",
     "• First inventory purchase\n• Hire 2 sales reps\n• Onboard 10 wholesalers\n• Activate M-Pesa payments\n• First 20 retailer orders"),
    (BLUE, "MONTH 2", "Scale & BNPL",
     "• KYC first cohort of retailers\n• Launch BNPL to verified retailers\n• First international product imports\n• 50+ active retailers\n• Track default rates"),
    (ORANGE, "MONTH 3", "Optimize & Grow",
     "• Expand product catalog\n• Add delivery logistics\n• 100+ active retailers\n• Revenue run-rate data\n• Prepare for seed round"),
]
for i, (color, label, title, items) in enumerate(months):
    x = Inches(0.8 + i * 4.0)
    y = Inches(2.5)
    add_rounded_rect(slide, x, y, Inches(3.6), Inches(4.0), GRAY_50, GRAY_200)
    # Top border color
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Inches(3.6), Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = color
    top_bar.line.fill.background()
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(3.0), Inches(0.3),
                 label, font_size=12, color=color, bold=True)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.7), Inches(3.0), Inches(0.4),
                 title, font_size=18, color=GRAY_900, bold=True)
    add_multiline_box(slide, x + Inches(0.3), y + Inches(1.3), Inches(3.0), Inches(2.5),
                      items.split("\n"), font_size=14, color=GRAY_700, line_spacing=1.5)

add_slide_number(slide, 13, TOTAL_SLIDES)
add_brand_tag(slide)

# ══════════════════════════════════════════════
# SLIDE 14: CLOSING
# ══════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_gradient_bg(slide, GREEN, GREEN_DARK)

# Logo
logo = add_rounded_rect(slide, Inches(0.8), Inches(1.2), Inches(1.1), Inches(1.1),
                         RGBColor(0xFF, 0xFF, 0xFF))
add_text_box(slide, Inches(0.8), Inches(1.18), Inches(1.1), Inches(1.1),
             "R", font_size=48, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)

add_text_box(slide, Inches(0.8), Inches(2.7), Inches(10), Inches(1.2),
             "Stocking them up.\nBacking them up.", font_size=48, color=WHITE, bold=True)

add_text_box(slide, Inches(0.8), Inches(4.3), Inches(8), Inches(0.8),
             "The product is built. The model works. The market is massive.\nWhat we need is capital to move inventory.",
             font_size=20, color=RGBColor(0xDC, 0xFC, 0xE7))

# CTA box
add_rounded_rect(slide, Inches(0.8), Inches(5.5), Inches(3.0), Inches(1.2),
                 RGBColor(0x12, 0x6B, 0x32), RGBColor(0x1A, 0x8B, 0x45))
add_text_box(slide, Inches(1.1), Inches(5.65), Inches(2.4), Inches(0.4),
             "Let's talk.", font_size=20, color=WHITE, bold=True)
add_text_box(slide, Inches(1.1), Inches(6.1), Inches(2.4), Inches(0.3),
             "ristoka.com", font_size=14, color=RGBColor(0xBB, 0xF7, 0xD0))

add_slide_number(slide, 14, TOTAL_SLIDES, RGBColor(0xBB, 0xF7, 0xD0))
add_brand_tag(slide, RGBColor(0xBB, 0xF7, 0xD0))

# ── Save ──
output_path = "public/ristoka-pitch-deck.pptx"
prs.save(output_path)
print(f"✓ Created {output_path} ({TOTAL_SLIDES} slides)")
